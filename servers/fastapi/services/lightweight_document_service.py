"""
Dependency-light native extraction for office documents.

The primary office path is LibreOffice -> PDF -> LiteParse
(``services/document_conversion_service.py``). That requires the ``soffice``
binary, which is absent on plain Windows installs of this fork. This service
extracts text directly for the modern (zip-based) formats we can handle without
any external binary, so office uploads still work when LibreOffice is missing.

``documents_loader`` wires this in as ``document_service`` (the optional fallback
converter). Formats it cannot handle natively (legacy OLE ``.doc``/``.ppt``/
``.xls`` and OpenDocument files) raise :class:`DocumentConversionError` with a
clear message so the caller can surface an actionable error instead of a raw
500. Modern OOXML spreadsheets (``.xlsx``/``.xlsm``) are parsed directly.
"""

import html
import posixpath
import re
import zipfile
from pathlib import Path
from xml.etree import ElementTree as ET

from services.document_conversion_service import DocumentConversionError

# python-pptx officially supports these; other extensions (.ppsx, legacy .ppt)
# are left to LibreOffice.
PPTX_EXTENSIONS = {".pptx", ".pptm"}
DOCX_EXTENSIONS = {".docx", ".docm"}
XLSX_EXTENSIONS = {".xlsx", ".xlsm"}
TEXT_EXTENSIONS = {".txt", ".csv", ".tsv"}


class DocumentService:
    def parse_to_markdown(self, file_path: str) -> str:
        ext = Path(file_path).suffix.lower()
        if ext in PPTX_EXTENSIONS:
            return _parse_pptx(file_path)
        if ext in DOCX_EXTENSIONS:
            return _parse_docx(file_path)
        if ext in XLSX_EXTENSIONS:
            return _parse_xlsx(file_path)
        if ext in TEXT_EXTENSIONS:
            return _read_text(file_path)
        raise DocumentConversionError(
            f"'{Path(file_path).name}' 형식은 LibreOffice 없이는 변환할 수 없습니다. "
            "PDF로 변환해 업로드하거나 LibreOffice를 설치해 주세요."
        )


def _parse_pptx(file_path: str) -> str:
    try:
        from pptx import Presentation

        prs = Presentation(file_path)
        blocks: list[str] = []
        for index, slide in enumerate(prs.slides, start=1):
            lines: list[str] = []
            for shape in slide.shapes:
                lines.extend(_shape_text_lines(shape))
            if slide.has_notes_slide:
                notes = slide.notes_slide.notes_text_frame
                if notes is not None and notes.text.strip():
                    lines.append(f"> {notes.text.strip()}")
            if lines:
                blocks.append(f"## Slide {index}\n" + "\n".join(lines))
        return "\n\n".join(blocks).strip()
    except DocumentConversionError:
        raise
    except Exception as exc:
        raise DocumentConversionError(
            f"'{Path(file_path).name}' PowerPoint 파일을 읽을 수 없습니다: {exc}"
        ) from exc


def _shape_text_lines(shape) -> list[str]:
    out: list[str] = []
    # Grouped shapes hold nested shapes (common in real decks) — recurse.
    if hasattr(shape, "shapes"):
        for child in shape.shapes:
            out.extend(_shape_text_lines(child))
        return out
    if getattr(shape, "has_text_frame", False):
        for para in shape.text_frame.paragraphs:
            text = "".join(run.text for run in para.runs).strip()
            if text:
                out.append(text)
    if getattr(shape, "has_table", False):
        for row in shape.table.rows:
            cells = [cell.text.strip() for cell in row.cells]
            if any(cells):
                out.append("| " + " | ".join(cells) + " |")
    return out


def _parse_docx(file_path: str) -> str:
    try:
        with zipfile.ZipFile(file_path) as zf:
            xml = zf.read("word/document.xml").decode("utf-8", errors="replace")
    except KeyError as exc:
        raise DocumentConversionError(
            f"'{Path(file_path).name}'에서 본문(word/document.xml)을 찾을 수 없습니다."
        ) from exc
    except Exception as exc:
        raise DocumentConversionError(
            f"'{Path(file_path).name}' Word 파일을 읽을 수 없습니다: {exc}"
        ) from exc

    # Turn structural markers into whitespace, then strip the remaining tags so
    # only the text nodes (the <w:t> contents) plus our breaks are left.
    xml = re.sub(r"</w:p>", "\n", xml)
    xml = re.sub(r"<w:br\b[^>]*/?>", "\n", xml)
    xml = re.sub(r"<w:tab\b[^>]*/?>", "\t", xml)
    text = re.sub(r"<[^>]+>", "", xml)
    text = html.unescape(text)
    text = re.sub(r"[ \t]+\n", "\n", text)
    text = re.sub(r"\n{3,}", "\n\n", text)
    return text.strip()


def _read_text(file_path: str) -> str:
    with open(file_path, "r", encoding="utf-8", errors="replace") as file:
        return file.read()


def _local_name(tag: str) -> str:
    return tag.rsplit("}", 1)[-1]


def _attribute_by_local_name(element: ET.Element, name: str) -> str | None:
    for key, value in element.attrib.items():
        if _local_name(key) == name:
            return value
    return None


def _read_xml(zf: zipfile.ZipFile, part_name: str) -> ET.Element:
    try:
        return ET.fromstring(zf.read(part_name))
    except KeyError as exc:
        raise DocumentConversionError(
            f"스프레드시트에 필요한 XML 파일({part_name})을 찾을 수 없습니다."
        ) from exc
    except ET.ParseError as exc:
        raise DocumentConversionError(
            f"스프레드시트 XML 파일({part_name})이 손상되었습니다: {exc}"
        ) from exc


def _shared_strings(zf: zipfile.ZipFile) -> list[str]:
    if "xl/sharedStrings.xml" not in zf.namelist():
        return []
    root = _read_xml(zf, "xl/sharedStrings.xml")
    values: list[str] = []
    for item in root.iter():
        if _local_name(item.tag) != "si":
            continue
        values.append(
            "".join(
                node.text or ""
                for node in item.iter()
                if _local_name(node.tag) == "t"
            )
        )
    return values


def _worksheet_parts(zf: zipfile.ZipFile) -> list[tuple[str, str]]:
    workbook = _read_xml(zf, "xl/workbook.xml")
    relationships: dict[str, str] = {}
    rels_name = "xl/_rels/workbook.xml.rels"
    if rels_name in zf.namelist():
        rels = _read_xml(zf, rels_name)
        for rel in rels.iter():
            if _local_name(rel.tag) != "Relationship":
                continue
            rel_id = rel.attrib.get("Id")
            target = rel.attrib.get("Target")
            rel_type = rel.attrib.get("Type", "")
            if rel_id and target and rel_type.endswith("/worksheet"):
                relationships[rel_id] = target

    sheets: list[tuple[str, str]] = []
    for sheet in workbook.iter():
        if _local_name(sheet.tag) != "sheet":
            continue
        name = sheet.attrib.get("name", f"Sheet {len(sheets) + 1}")
        rel_id = _attribute_by_local_name(sheet, "id")
        target = relationships.get(rel_id or "")
        if not target:
            continue
        part_name = (
            target.lstrip("/")
            if target.startswith("/")
            else posixpath.normpath(posixpath.join("xl", target))
        )
        if part_name in zf.namelist():
            sheets.append((name, part_name))

    if sheets:
        return sheets

    # Some producers omit relationship metadata. Keep a useful best-effort
    # fallback instead of rejecting otherwise readable worksheet XML.
    fallback_parts = sorted(
        name
        for name in zf.namelist()
        if name.startswith("xl/worksheets/") and name.endswith(".xml")
    )
    return [(f"Sheet {index}", name) for index, name in enumerate(fallback_parts, 1)]


def _column_index(cell_reference: str) -> int:
    letters = re.match(r"[A-Za-z]+", cell_reference or "")
    if not letters:
        return 0
    index = 0
    for char in letters.group(0).upper():
        index = index * 26 + (ord(char) - ord("A") + 1)
    return max(index - 1, 0)


def _cell_value(cell: ET.Element, shared: list[str]) -> str:
    cell_type = cell.attrib.get("t", "")
    value_node = next(
        (node for node in cell if _local_name(node.tag) == "v"),
        None,
    )
    raw_value = value_node.text if value_node is not None and value_node.text else ""

    if cell_type == "inlineStr":
        return "".join(
            node.text or ""
            for node in cell.iter()
            if _local_name(node.tag) == "t"
        )
    if cell_type == "s":
        try:
            return shared[int(raw_value)]
        except (ValueError, IndexError):
            return raw_value
    if cell_type == "b":
        return "TRUE" if raw_value == "1" else "FALSE"
    if raw_value:
        return raw_value

    formula = next(
        (node.text for node in cell if _local_name(node.tag) == "f" and node.text),
        None,
    )
    return f"={formula}" if formula else ""


def _markdown_cell(value: str) -> str:
    return (
        value.replace("\r\n", " ")
        .replace("\r", " ")
        .replace("\n", " ")
        .replace("|", "\\|")
        .strip()
    )


def _parse_worksheet(
    zf: zipfile.ZipFile,
    part_name: str,
    shared: list[str],
) -> list[str]:
    root = _read_xml(zf, part_name)
    output: list[str] = []
    for row in root.iter():
        if _local_name(row.tag) != "row":
            continue
        values: list[str] = []
        next_column = 0
        for cell in row:
            if _local_name(cell.tag) != "c":
                continue
            column = _column_index(cell.attrib.get("r", ""))
            if "r" not in cell.attrib:
                column = next_column
            if column > 16_383:  # XLSX's maximum column (XFD).
                continue
            if column >= len(values):
                values.extend([""] * (column - len(values) + 1))
            values[column] = _markdown_cell(_cell_value(cell, shared))
            next_column = column + 1
        while values and not values[-1]:
            values.pop()
        if values:
            output.append("| " + " | ".join(values) + " |")
    return output


def _parse_xlsx(file_path: str) -> str:
    try:
        with zipfile.ZipFile(file_path) as zf:
            shared = _shared_strings(zf)
            blocks: list[str] = []
            for sheet_name, part_name in _worksheet_parts(zf):
                rows = _parse_worksheet(zf, part_name, shared)
                if rows:
                    blocks.append(f"## {sheet_name}\n" + "\n".join(rows))
            return "\n\n".join(blocks).strip()
    except DocumentConversionError:
        raise
    except zipfile.BadZipFile as exc:
        raise DocumentConversionError(
            f"'{Path(file_path).name}' Excel 파일이 손상되었거나 올바른 XLSX 형식이 아닙니다."
        ) from exc
    except Exception as exc:
        raise DocumentConversionError(
            f"'{Path(file_path).name}' Excel 파일을 읽을 수 없습니다: {exc}"
        ) from exc
