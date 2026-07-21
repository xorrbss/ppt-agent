"""Offline PPTX design-signal extraction with OOXML safety preflight."""

from __future__ import annotations

from collections import Counter, defaultdict
from hashlib import sha256
from io import BytesIO
from pathlib import Path
from typing import Any, Iterable
from xml.etree import ElementTree
from zipfile import BadZipFile, ZipFile, is_zipfile
import zlib

from pptx import Presentation
from pptx.enum.dml import MSO_COLOR_TYPE
from pptx.enum.shapes import MSO_SHAPE_TYPE, PP_PLACEHOLDER
from pptx.exc import PackageNotFoundError

from utils.artifact_style_analysis import (
    MAX_ARTIFACT_BYTES,
    MAX_PAGES,
    ArtifactAnalysisError,
)


MAX_PACKAGE_MEMBERS = 10_000
MAX_UNCOMPRESSED_BYTES = 4 * MAX_ARTIFACT_BYTES
_DRAWING_NS = "http://schemas.openxmlformats.org/drawingml/2006/main"
_PACKAGE_MEMBER_ERRORS = (BadZipFile, NotImplementedError, RuntimeError, zlib.error)


def _warning(code: str, message: str) -> dict[str, str]:
    return {"code": code, "message": message}


def _local_name(tag: str) -> str:
    return tag.rsplit("}", 1)[-1]


def _inspect_package(
    raw: bytes,
) -> tuple[dict[str, Any], list[dict[str, str]], ZipFile]:
    if not is_zipfile(BytesIO(raw)):
        raise ArtifactAnalysisError("PPTX input is not a valid ZIP-based OOXML package")
    try:
        package = ZipFile(BytesIO(raw))
        members = package.infolist()
    except BadZipFile as exc:
        raise ArtifactAnalysisError(f"PPTX package is damaged: {exc}") from exc
    if len(members) > MAX_PACKAGE_MEMBERS:
        package.close()
        raise ArtifactAnalysisError(
            f"PPTX package exceeds the {MAX_PACKAGE_MEMBERS}-member safety limit"
        )
    if any(member.flag_bits & 0x1 for member in members):
        package.close()
        raise ArtifactAnalysisError(
            "encrypted ZIP members are not supported in PPTX input"
        )
    uncompressed = sum(member.file_size for member in members)
    if uncompressed > MAX_UNCOMPRESSED_BYTES:
        package.close()
        raise ArtifactAnalysisError(
            "PPTX expanded content exceeds the safe uncompressed-size limit"
        )

    names = [member.filename for member in members]
    required = {"[Content_Types].xml", "ppt/presentation.xml"}
    if not required.issubset(names):
        package.close()
        raise ArtifactAnalysisError("PPTX package is missing required OOXML parts")
    for member in members:
        parts = member.filename.replace("\\", "/").split("/")
        if member.filename.startswith(("/", "\\")) or ".." in parts:
            package.close()
            raise ArtifactAnalysisError("PPTX package contains an unsafe member path")
        if member.file_size > MAX_ARTIFACT_BYTES or (
            member.compress_size > 0
            and member.file_size > 10 * 1024 * 1024
            and member.file_size / member.compress_size > 1000
        ):
            package.close()
            raise ArtifactAnalysisError(
                "PPTX package contains an unsafe compressed member"
            )
    macro_parts = sorted(
        name for name in names if name.lower().endswith("vbaproject.bin")
    )
    embedded_parts = sorted(
        name for name in names if name.startswith("ppt/embeddings/")
    )
    active_parts = sorted(
        name
        for name in names
        if name.startswith("ppt/activeX/") or name.startswith("ppt/ctrlProps/")
    )
    external_links = 0
    for name in sorted(name for name in names if name.endswith(".rels")):
        try:
            root = ElementTree.fromstring(package.read(name))
        except _PACKAGE_MEMBER_ERRORS as exc:
            package.close()
            raise ArtifactAnalysisError(
                f"PPTX package contains a damaged or unsupported member: {exc}"
            ) from exc
        except (ElementTree.ParseError, KeyError):
            continue
        external_links += sum(
            1
            for node in root
            if node.attrib.get("TargetMode", "").casefold() == "external"
        )

    warnings: list[dict[str, str]] = []
    if macro_parts:
        warnings.append(
            _warning(
                "pptx_macros_present",
                "Macro parts were detected and were not executed.",
            )
        )
    if external_links:
        warnings.append(
            _warning(
                "pptx_external_links_present",
                f"{external_links} external relationship(s) were detected and not followed.",
            )
        )
    if embedded_parts:
        warnings.append(
            _warning(
                "pptx_embedded_content_present",
                f"{len(embedded_parts)} embedded part(s) were detected and not opened.",
            )
        )
    if active_parts:
        warnings.append(
            _warning(
                "pptx_active_content_present",
                f"{len(active_parts)} ActiveX/control part(s) were detected and not opened.",
            )
        )
    security = {
        "active_content_part_count": len(active_parts),
        "embedded_part_count": len(embedded_parts),
        "external_relationship_count": external_links,
        "macro_part_count": len(macro_parts),
        "policy": "Macros, external links, and embedded content were not executed.",
    }
    return security, warnings, package


def _theme_values(package: ZipFile) -> tuple[list[str], list[str]]:
    colors: list[str] = []
    fonts: list[str] = []
    names = sorted(
        name
        for name in package.namelist()
        if name.startswith("ppt/theme/") and name.endswith(".xml")
    )
    for name in names:
        try:
            root = ElementTree.fromstring(package.read(name))
        except (ElementTree.ParseError, KeyError):
            continue
        for node in root.findall(f".//{{{_DRAWING_NS}}}clrScheme/*"):
            for color_node in node:
                value = color_node.attrib.get("val") or color_node.attrib.get("lastClr")
                if value and len(value) == 6:
                    candidate = f"#{value.upper()}"
                    if candidate not in colors:
                        colors.append(candidate)
                    break
        for node in root.findall(f".//{{{_DRAWING_NS}}}fontScheme//*"):
            if _local_name(node.tag) in {"latin", "ea", "cs"}:
                value = node.attrib.get("typeface", "").strip()
                if value and value not in fonts:
                    fonts.append(value)
    return colors, fonts


def _walk_shapes(shapes: Iterable[Any]) -> Iterable[Any]:
    for shape in shapes:
        if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
            yield from _walk_shapes(shape.shapes)
        else:
            yield shape


def _shape_kind(shape: Any) -> str:
    if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
        return "image"
    if getattr(shape, "has_table", False):
        return "table"
    if getattr(shape, "has_chart", False):
        return "chart"
    if getattr(shape, "has_text_frame", False) and shape.text.strip():
        return "text"
    return "shape"


def _occupancy_cells(
    shape: Any, width: int, height: int, *, grid_size: int = 100
) -> set[int]:
    """Return deterministic canvas cells covered by a top-level shape box."""
    if width <= 0 or height <= 0:
        return set()
    left_value = min(width, max(0, int(shape.left)))
    top_value = min(height, max(0, int(shape.top)))
    right_value = min(width, max(left_value, left_value + max(0, int(shape.width))))
    bottom_value = min(height, max(top_value, top_value + max(0, int(shape.height))))
    if right_value <= left_value or bottom_value <= top_value:
        return set()
    left = min(grid_size - 1, int(left_value / width * grid_size))
    right = min(
        grid_size,
        max(left + 1, (right_value * grid_size + width - 1) // width),
    )
    first_row = min(grid_size - 1, int(top_value / height * grid_size))
    last_row = min(
        grid_size,
        max(first_row + 1, (bottom_value * grid_size + height - 1) // height),
    )
    return {
        row * grid_size + column
        for row in range(first_row, last_row)
        for column in range(left, right)
    }


def _add_color(
    color_format: Any, role: str, colors: Counter[str], roles: dict[str, set[str]]
) -> None:
    try:
        if color_format.type == MSO_COLOR_TYPE.RGB and color_format.rgb is not None:
            value = f"#{str(color_format.rgb).upper()}"
            colors[value] += 1
            roles[value].add(role)
    except (AttributeError, TypeError, ValueError):
        return


def _placeholder_role(shape: Any, level: int) -> str:
    try:
        placeholder_type = shape.placeholder_format.type
    except (AttributeError, ValueError):
        placeholder_type = None
    if placeholder_type in {PP_PLACEHOLDER.TITLE, PP_PLACEHOLDER.CENTER_TITLE}:
        return "title"
    if placeholder_type == PP_PLACEHOLDER.SUBTITLE:
        return "subtitle"
    return "body" if level == 0 else f"body-level-{level}"


def _layout_signature(
    shapes: list[Any], width: int, height: int, layout_name: str
) -> str:
    cells: Counter[str] = Counter()
    for shape in shapes:
        kind = _shape_kind(shape)
        x = min(3, max(0, int((max(0, shape.left) / width) * 4)))
        y = min(3, max(0, int((max(0, shape.top) / height) * 4)))
        w = min(4, max(1, round((max(0, shape.width) / width) * 4)))
        h = min(4, max(1, round((max(0, shape.height) / height) * 4)))
        cells[f"{kind}:{x},{y},{w},{h}"] += 1
    payload = f"{layout_name}|" + "|".join(
        f"{key}={min(value, 9)}" for key, value in sorted(cells.items())
    )
    return sha256(payload.encode("utf-8")).hexdigest()[:12]


def _signal(values: list[dict[str, Any]], confidence: str) -> dict[str, Any]:
    return {
        "confidence": confidence if values else "none",
        "status": "observed" if values else "unavailable",
        "values": values,
    }


def analyze_pptx(source: Path, raw: bytes) -> dict[str, Any]:
    security, warnings, package = _inspect_package(raw)
    try:
        theme_colors, theme_fonts = _theme_values(package)
    except _PACKAGE_MEMBER_ERRORS as exc:
        raise ArtifactAnalysisError(
            f"{source}: PPTX package contains a damaged or unsupported member: {exc}"
        ) from exc
    finally:
        package.close()
    try:
        presentation = Presentation(BytesIO(raw))
    except (
        PackageNotFoundError,
        BadZipFile,
        KeyError,
        ValueError,
        NotImplementedError,
        RuntimeError,
        zlib.error,
    ) as exc:
        raise ArtifactAnalysisError(
            f"{source}: damaged or unsupported PPTX: {exc}"
        ) from exc
    slide_count = len(presentation.slides)
    if slide_count == 0:
        raise ArtifactAnalysisError(f"{source}: PPTX contains no slides")
    if slide_count > MAX_PAGES:
        raise ArtifactAnalysisError(
            f"{source}: PPTX exceeds the {MAX_PAGES}-slide safety limit"
        )

    # A malformed presentation.xml can omit <p:sldSz> (slide_width None) or declare
    # a zero dimension; validate here — mirroring the PDF page-dimension guard — so
    # the aspect-ratio math below can't raise a TypeError/ZeroDivisionError that
    # would escape the CLI's ArtifactAnalysisError handler as a raw traceback.
    raw_width = presentation.slide_width
    raw_height = presentation.slide_height
    if raw_width is None or raw_height is None:
        raise ArtifactAnalysisError(f"{source}: PPTX has no slide size")
    width = int(raw_width)
    height = int(raw_height)
    if width <= 0 or height <= 0:
        raise ArtifactAnalysisError(f"{source}: PPTX has invalid slide dimensions")
    colors: Counter[str] = Counter()
    color_roles: dict[str, set[str]] = defaultdict(set)
    fonts: Counter[str] = Counter()
    font_sources: dict[str, set[str]] = defaultdict(set)
    sizes: Counter[tuple[float, str]] = Counter()
    kind_counts: Counter[str] = Counter()
    occupied_cells: dict[str, set[tuple[int, int]]] = defaultdict(set)
    signatures: dict[str, list[int]] = defaultdict(list)
    layout_names: dict[str, Counter[str]] = defaultdict(Counter)

    for color in theme_colors:
        colors[color] += 1
        color_roles[color].add("theme")
    for font in theme_fonts:
        fonts[font] += 1
        font_sources[font].add("theme")

    for slide_number, slide in enumerate(presentation.slides, start=1):
        try:
            fill = slide.background.fill
            if fill.type is not None:
                _add_color(fill.fore_color, "background", colors, color_roles)
        except (AttributeError, TypeError, ValueError):
            pass
        top_level_shapes = list(slide.shapes)
        shapes = list(_walk_shapes(top_level_shapes))
        layout_name = (slide.slide_layout.name or "unnamed").strip()
        signature = _layout_signature(top_level_shapes, width, height, layout_name)
        signatures[signature].append(slide_number)
        layout_names[signature][layout_name] += 1

        for shape in top_level_shapes:
            kind = _shape_kind(shape)
            occupied_cells[kind].update(
                (slide_number, cell) for cell in _occupancy_cells(shape, width, height)
            )

        for shape in shapes:
            kind = _shape_kind(shape)
            kind_counts[kind] += 1
            try:
                if shape.fill.type is not None:
                    _add_color(shape.fill.fore_color, "fill", colors, color_roles)
                _add_color(shape.line.color, "line", colors, color_roles)
            except (AttributeError, TypeError, ValueError):
                pass
            if not getattr(shape, "has_text_frame", False):
                continue
            for paragraph in shape.text_frame.paragraphs:
                role = _placeholder_role(shape, paragraph.level)
                for run in paragraph.runs:
                    if run.font.name and run.font.name.strip():
                        family = run.font.name.strip()
                        fonts[family] += max(1, len(run.text.strip()))
                        font_sources[family].add("run")
                    if run.font.size is not None:
                        sizes[(round(run.font.size.pt, 1), role)] += max(
                            1, len(run.text.strip())
                        )
                    _add_color(run.font.color, "text", colors, color_roles)

    color_values = [
        {"count": count, "roles": sorted(color_roles[value]), "value": value}
        for value, count in sorted(
            colors.items(), key=lambda item: (-item[1], item[0])
        )[:12]
    ]
    font_values = [
        {"count": count, "family": family, "sources": sorted(font_sources[family])}
        for family, count in sorted(
            fonts.items(), key=lambda item: (-item[1], item[0].casefold(), item[0])
        )[:12]
    ]
    hierarchy_values = [
        {"count": count, "role": role, "size_pt": size}
        for (size, role), count in sorted(
            sizes.items(), key=lambda item: (-item[0][0], item[0][1], -item[1])
        )[:12]
    ]
    repeated_values = []
    for signature, slides in sorted(
        signatures.items(), key=lambda item: (-len(item[1]), item[0])
    ):
        if len(slides) < 2:
            continue
        name = sorted(
            layout_names[signature].items(),
            key=lambda item: (-item[1], item[0].casefold(), item[0]),
        )[0][0]
        repeated_values.append(
            {
                "count": len(slides),
                "layout_name": name,
                "signature": signature,
                "slides": slides,
            }
        )

    total_occupied_cells = sum(len(cells) for cells in occupied_cells.values())
    area_share = {
        kind: round(len(cells) / total_occupied_cells, 4)
        for kind, cells in sorted(occupied_cells.items())
        if total_occupied_cells
    }
    canvas_coverage = {
        kind: round(len(cells) / (slide_count * 10_000), 4)
        for kind, cells in sorted(occupied_cells.items())
    }
    if not color_values:
        warnings.append(
            _warning(
                "colors_unavailable",
                "No explicit RGB or theme colors were extractable.",
            )
        )
    if not font_values:
        warnings.append(
            _warning(
                "fonts_unavailable",
                "No explicit run or theme font families were extractable.",
            )
        )
    if not hierarchy_values:
        warnings.append(
            _warning(
                "text_hierarchy_unavailable", "No explicit text sizes were extractable."
            )
        )
    if not repeated_values:
        warnings.append(
            _warning(
                "repeated_layouts_unavailable",
                "No repeated coarse slide layout was observed.",
            )
        )

    return {
        "document": {
            "page_count": slide_count,
            "page_size": {
                "aspect_ratio": round(width / height, 6),
                "height": round(height / 914400, 4),
                "unit": "inches",
                "width": round(width / 914400, 4),
            },
        },
        "evidence": [
            {
                "method": "OOXML theme and DrawingML RGB values",
                "observations": len(color_values),
                "signal": "colors",
            },
            {
                "method": "theme font scheme and explicit text runs",
                "observations": len(font_values),
                "signal": "fonts",
            },
            {
                "method": "explicit run sizes and placeholder roles",
                "observations": sum(sizes.values()),
                "signal": "text_hierarchy",
            },
            {
                "method": "coarse normalized shape geometry",
                "observations": len(repeated_values),
                "signal": "repeated_layouts",
            },
        ],
        "security": security,
        "signals": {
            "colors": _signal(
                color_values, "high" if len(color_values) >= 3 else "medium"
            ),
            "composition": {
                "canvas_coverage": canvas_coverage,
                "confidence": "medium" if kind_counts else "none",
                "element_area_share": area_share,
                "element_counts": dict(sorted(kind_counts.items())),
                "note": "Shares use deterministic 100x100 per-kind canvas occupancy; cross-kind overlaps remain attributable to each kind.",
                "status": "observed" if kind_counts else "unavailable",
            },
            "fonts": _signal(
                font_values,
                "high"
                if any("run" in item["sources"] for item in font_values)
                else "medium",
            ),
            "repeated_layouts": _signal(repeated_values, "medium"),
            "text_hierarchy": _signal(
                hierarchy_values, "high" if hierarchy_values else "none"
            ),
        },
        "warnings": warnings,
    }
