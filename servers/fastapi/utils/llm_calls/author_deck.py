"""Authored-mode deck builder: turn an outline into a list of bespoke HTML slides
(authored concurrently against one shared design system for cohesion) and assemble
the rendered slide images into an image-per-slide PPTX that opens in PowerPoint with
perfect fidelity. The render step itself reuses utils/slide_capture. Opt-in — the
fast default template path is untouched."""

import asyncio
import io
import logging
from dataclasses import dataclass
from typing import List, Optional

from models.presentation_outline_model import PresentationOutlineModel
from utils.llm_calls.author_slide import (
    AuthoredStyleLike,
    Brand,
    author_slide_html,
    apply_style_defaults,
    build_design_system,
    fallback_slide_html,
    is_valid_slide_html,
)
from utils.authored_styles import reference_images_for_role

logger = logging.getLogger(__name__)


@dataclass(frozen=True)
class AuthoredDeckResult:
    """First-pass slides plus the exact design system reused by vision-QA."""

    htmls: List[str]
    design_system: str


# Concurrency is bounded PROCESS-WIDE inside author_slide_html (the AUTHOR_CONCURRENCY-sized
# dedicated executor), so the per-slide fan-out here doesn't need its own semaphore — excess
# calls queue at that single chokepoint shared with the vision-QA re-author pass.

# A light deck-plan: a slide's ROLE nudges the model toward a fitting bespoke layout
# (cover hero / editorial problem / numbered pillars / phased roadmap / hero metric /
# bold closing). Position drives the frame; content keywords pick the middle role.
_METRIC_HINTS = (
    "%", "퍼센트", "성장", "매출", "수익", "지표", "비율", "증가", "감소", "달성",
    "ROI", "조원", "억원", "billion", "million", "growth", "revenue",
)
_TIMELINE_HINTS = (
    "단계", "로드맵", "분기", "타임라인", "절차", "순서", "스텝",
    "단계", "로드맵", "프로세스", "대응 흐름", "phase", "roadmap", "timeline", "step",
    "pipeline", "quarter", "2024", "2025", "2026", "2027",
)
_ARCHITECTURE_HINTS = (
    "아키텍처", "구조", "시스템", "에이전트", "architecture", "system", "platform",
    "layer", "integration", "agent",
)


def derive_role(index: int, n: int, content: str) -> str:
    """Per-slide ROLE for the authoring prompt, from position + content keywords."""
    if index == 0:
        return "COVER"
    if index == n - 1:
        return "CLOSING"
    if index == 1:
        return "PROBLEM"
    lowered = content.lower()
    if any(h.lower() in lowered for h in _ARCHITECTURE_HINTS):
        return "ARCHITECTURE"
    if any(h.lower() in lowered for h in _TIMELINE_HINTS):
        return "ROADMAP"
    if any(h.lower() in lowered for h in _METRIC_HINTS):
        return "OUTCOMES"
    return "PILLARS"


def plan_deck_roles(outline: PresentationOutlineModel) -> List[str]:
    """The deck-plan: one ROLE per slide. Single source of role derivation, shared by
    the authoring pass and the vision-QA re-author pass."""
    slides = list(outline.slides)
    n = len(slides)
    return [derive_role(i, n, slides[i].content) for i in range(n)]


async def _author_one_resilient(
    content: str,
    design_system: str,
    brand: Brand,
    role: str,
    index: int,
    n: int,
    reference_images=(),
) -> str:
    """Author one slide with one retry and a branded fallback. Never raises and never
    returns invalid HTML — so one bad slide degrades to a clean placeholder instead of
    aborting the whole deck. (Concurrency is bounded inside author_slide_html.) A
    fallback is logged at WARNING so a silently-degraded deck is visible to operators."""
    last_error: object = "invalid HTML"
    for _ in range(2):
        try:
            if reference_images:
                html = await author_slide_html(
                    content,
                    design_system,
                    brand,
                    role,
                    index,
                    n,
                    reference_images=reference_images,
                )
            else:
                html = await author_slide_html(
                    content, design_system, brand, role, index, n
                )
            if is_valid_slide_html(html):
                return html
        except Exception as e:  # noqa: BLE001 — resilience boundary, recorded below
            last_error = e
    logger.warning(
        "[authored] slide %d/%d fell back to a branded placeholder after 2 attempts "
        "(last error: %s)",
        index + 1,
        n,
        last_error,
    )
    return fallback_slide_html(content, brand, role, index, n)


async def author_deck(
    outline: PresentationOutlineModel,
    brand: Brand,
    style: Optional[AuthoredStyleLike] = None,
) -> AuthoredDeckResult:
    """Author every outline slide against the shared design system, with bounded
    concurrency and per-slide resilience. Returns the valid HTML documents (order
    preserved) plus the exact prompt shared with vision-QA; failed slides fall back
    to a clean branded placeholder."""
    slides = list(outline.slides)
    n = len(slides)
    brand = apply_style_defaults(brand, style)
    design_system = build_design_system(brand, style)
    roles = plan_deck_roles(outline)
    htmls = await asyncio.gather(
        *[
            _author_one_resilient(
                slides[i].content,
                design_system,
                brand,
                roles[i],
                i,
                n,
                reference_images_for_role(style, roles[i]),
            )
            for i in range(n)
        ]
    )
    return AuthoredDeckResult(htmls=list(htmls), design_system=design_system)


def build_image_pptx(images: List[bytes], out_path: str) -> str:
    """Assemble rendered slide PNGs into a 16:9 PPTX, one full-bleed picture per
    slide (perfect fidelity, opens in PowerPoint). Returns the saved path."""
    from pptx import Presentation
    from pptx.util import Inches

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]
    for img in images:
        slide = prs.slides.add_slide(blank)
        slide.shapes.add_picture(
            io.BytesIO(img), 0, 0, width=prs.slide_width, height=prs.slide_height
        )
    prs.save(out_path)
    return out_path
