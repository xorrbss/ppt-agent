from __future__ import annotations

from io import BytesIO
import json
from pathlib import Path
from zipfile import ZipFile

import pytest
from pptx import Presentation

from templates.v2.models.layouts import SlideLayout
from templates.v2.pptx.native_compiler import (
    TEMPLATE_V2_NATIVE_PPTX_COMPILER_NAME,
    TEMPLATE_V2_NATIVE_PPTX_COMPILER_VERSION,
    TEMPLATE_V2_NATIVE_PPTX_SCHEMA_VERSION,
    TemplateV2NativePptxCompileError,
    compile_template_v2_pptx,
    inspect_native_pptx_structure,
)

FIXTURE = (
    Path(__file__).parents[1]
    / "fixtures"
    / "template_v2"
    / "native_pptx_golden_v1.json"
)


def _golden_layout() -> SlideLayout:
    return SlideLayout.model_validate(
        {
            "id": "native-golden",
            "description": "Deterministic native compiler golden layout",
            "components": [
                {
                    "id": "hero",
                    "description": "Editable hero and decorative background",
                    "position": {"x": 0, "y": 0},
                    "elements": [
                        {
                            "type": "container",
                            "position": {"x": 0, "y": 0},
                            "size": {"width": 1280, "height": 720},
                            "fill": {"color": "#F4F7FB"},
                            "stroke": {"color": "#D8E1EE", "width": 1},
                        },
                        {
                            "type": "text",
                            "position": {"x": 96, "y": 88},
                            "size": {"width": 1088, "height": 90},
                            "runs": [{"text": "Original title"}],
                            "font": {
                                "size": 32,
                                "family": "Arial",
                                "color": "#123456",
                                "bold": True,
                            },
                            "alignment": {
                                "horizontal": "left",
                                "vertical": "middle",
                            },
                            "decorative": False,
                            "name": "title",
                            "min_length": 1,
                            "max_length": 80,
                        },
                        {
                            "type": "text",
                            "position": {"x": 96, "y": 194},
                            "size": {"width": 1088, "height": 55},
                            "runs": [{"text": "Native OOXML · editable text"}],
                            "font": {
                                "size": 18,
                                "family": "Arial",
                                "color": "#52657A",
                            },
                            "decorative": True,
                            "name": "eyebrow",
                            "min_length": 1,
                            "max_length": 80,
                        },
                        {
                            "type": "group",
                            "name": "details",
                            "position": {"x": 96, "y": 300},
                            "size": {"width": 1088, "height": 240},
                            "children": [
                                {
                                    "type": "text-list",
                                    "position": {"x": 0, "y": 0},
                                    "size": {"width": 1088, "height": 240},
                                    "font": {
                                        "size": 22,
                                        "family": "Arial",
                                        "color": "#123456",
                                    },
                                    "marker": "none",
                                    "items": [[{"text": "Original item"}]],
                                    "decorative": False,
                                    "name": "items",
                                    "min_items": 1,
                                    "max_items": 3,
                                    "min_item_length": 1,
                                    "max_item_length": 80,
                                }
                            ],
                        },
                    ],
                }
            ],
        }
    )


def _golden_content() -> dict:
    return {
        "hero": {
            "title": "Quarterly direction",
            "details": {"items": ["Plan", "Build", "Measure"]},
        }
    }


def test_native_compiler_is_byte_deterministic_and_matches_golden():
    first = compile_template_v2_pptx([(_golden_layout(), _golden_content())])
    second = compile_template_v2_pptx([(_golden_layout(), _golden_content())])

    assert first.pptx == second.pptx
    assert first.package_sha256 == second.package_sha256
    assert first.structural_sha256 == second.structural_sha256
    assert first.manifest == json.loads(FIXTURE.read_text(encoding="utf-8"))
    assert first.manifest["schema_version"] == TEMPLATE_V2_NATIVE_PPTX_SCHEMA_VERSION
    assert first.manifest["compiler"] == {
        "name": TEMPLATE_V2_NATIVE_PPTX_COMPILER_NAME,
        "version": TEMPLATE_V2_NATIVE_PPTX_COMPILER_VERSION,
    }

    with ZipFile(BytesIO(first.pptx)) as package:
        assert package.namelist() == sorted(package.namelist())
        assert {entry.date_time for entry in package.infolist()} == {
            (1980, 1, 1, 0, 0, 0)
        }


def test_native_pptx_keeps_text_editable_and_geometry_at_96_dpi():
    compilation = compile_template_v2_pptx(
        [(_golden_layout(), _golden_content())]
    )
    presentation = Presentation(BytesIO(compilation.pptx))
    slide = presentation.slides[0]

    assert presentation.core_properties.subject == (
        TEMPLATE_V2_NATIVE_PPTX_SCHEMA_VERSION
    )
    assert presentation.core_properties.author == (
        TEMPLATE_V2_NATIVE_PPTX_COMPILER_NAME
    )
    assert presentation.core_properties.comments == (
        f"{TEMPLATE_V2_NATIVE_PPTX_COMPILER_NAME}/"
        f"{TEMPLATE_V2_NATIVE_PPTX_COMPILER_VERSION}"
    )
    assert (presentation.slide_width, presentation.slide_height) == (
        12_192_000,
        6_858_000,
    )
    assert [shape.name for shape in slide.shapes] == [
        "hero/container-0",
        "hero/title",
        "hero/eyebrow",
        "hero/details/items",
    ]
    assert [shape.text for shape in slide.shapes if shape.has_text_frame] == [
        "",
        "Quarterly direction",
        "Native OOXML · editable text",
        "Plan\nBuild\nMeasure",
    ]
    title = slide.shapes[1]
    assert (title.left, title.top, title.width, title.height) == (
        914_400,
        838_200,
        10_363_200,
        857_250,
    )
    assert title.text_frame.paragraphs[0].runs[0].font.name == "Arial"
    assert title.text_frame.paragraphs[0].runs[0].font.size.pt == 32

    assert inspect_native_pptx_structure(
        compilation.pptx,
        layout_ids=["native-golden"],
    ) == compilation.manifest


def test_native_compiler_fails_closed_for_unsupported_element():
    layout = SlideLayout.model_validate(
        {
            "id": "chart",
            "description": "Unsupported chart compiler contract",
            "components": [
                {
                    "id": "chart-panel",
                    "description": "Editable chart component content",
                    "position": {"x": 0, "y": 0},
                    "elements": [
                        {
                            "type": "chart",
                            "position": {"x": 80, "y": 80},
                            "size": {"width": 640, "height": 360},
                            "chart_type": "bar",
                            "categories": ["A"],
                            "series": [{"name": "Series", "values": [1]}],
                            "decorative": False,
                            "name": "chart",
                        }
                    ],
                }
            ],
        }
    )
    content = {
        "chart-panel": {
            "chart": {
                "chart_type": "bar",
                "categories": ["A"],
                "series": [{"name": "Series", "values": [1]}],
            }
        }
    }

    with pytest.raises(
        TemplateV2NativePptxCompileError,
        match="template_v2_native_pptx_element_unsupported",
    ) as error:
        compile_template_v2_pptx([(layout, content)])

    assert error.value.path.endswith(".type=chart")
