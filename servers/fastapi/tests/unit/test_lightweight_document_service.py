import io
import os
import tempfile
import zipfile
from unittest.mock import MagicMock

import pytest

from services.document_conversion_service import DocumentConversionError
from services.documents_loader import DocumentsLoader
from services.lightweight_document_service import DocumentService


def _make_pptx(path: str) -> None:
    from pptx import Presentation
    from pptx.util import Inches

    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "제안서 제목"
    body = slide.placeholders[1].text_frame
    body.text = "첫 번째 항목"
    body.add_paragraph().text = "두 번째 항목"

    table_slide = prs.slides.add_slide(prs.slide_layouts[5])
    table = table_slide.shapes.add_table(
        2, 2, Inches(1), Inches(1), Inches(4), Inches(1)
    ).table
    table.cell(0, 0).text = "구분"
    table.cell(0, 1).text = "금액"
    prs.save(path)


def _make_docx(path: str, body_xml: str) -> None:
    document_xml = (
        '<?xml version="1.0"?>'
        '<w:document xmlns:w="x"><w:body>' + body_xml + "</w:body></w:document>"
    )
    with zipfile.ZipFile(path, "w") as zf:
        zf.writestr("word/document.xml", document_xml)


def _make_xlsx(path: str) -> None:
    workbook_xml = """<?xml version="1.0" encoding="UTF-8"?>
    <workbook xmlns="http://schemas.openxmlformats.org/spreadsheetml/2006/main"
      xmlns:r="http://schemas.openxmlformats.org/officeDocument/2006/relationships">
      <sheets><sheet name="요약" sheetId="1" r:id="rId1"/></sheets>
    </workbook>"""
    rels_xml = """<?xml version="1.0" encoding="UTF-8"?>
    <Relationships xmlns="http://schemas.openxmlformats.org/package/2006/relationships">
      <Relationship Id="rId1"
        Type="http://schemas.openxmlformats.org/officeDocument/2006/relationships/worksheet"
        Target="worksheets/sheet1.xml"/>
    </Relationships>"""
    shared_strings_xml = """<?xml version="1.0" encoding="UTF-8"?>
    <sst xmlns="http://schemas.openxmlformats.org/spreadsheetml/2006/main">
      <si><t>항목</t></si><si><t>매출</t></si>
    </sst>"""
    sheet_xml = """<?xml version="1.0" encoding="UTF-8"?>
    <worksheet xmlns="http://schemas.openxmlformats.org/spreadsheetml/2006/main">
      <sheetData>
        <row r="1"><c r="A1" t="s"><v>0</v></c><c r="B1" t="inlineStr"><is><t>금액</t></is></c></row>
        <row r="2"><c r="A2" t="s"><v>1</v></c><c r="B2"><v>1200</v></c><c r="C2" t="b"><v>1</v></c></row>
        <row r="3"><c r="B3"><f>SUM(B2:B2)</f></c></row>
      </sheetData>
    </worksheet>"""
    with zipfile.ZipFile(path, "w") as zf:
        zf.writestr("xl/workbook.xml", workbook_xml)
        zf.writestr("xl/_rels/workbook.xml.rels", rels_xml)
        zf.writestr("xl/sharedStrings.xml", shared_strings_xml)
        zf.writestr("xl/worksheets/sheet1.xml", sheet_xml)


def test_parse_pptx_extracts_text_and_tables(tmp_path):
    path = str(tmp_path / "deck.pptx")
    _make_pptx(path)

    out = DocumentService().parse_to_markdown(path)

    assert "제안서 제목" in out
    assert "첫 번째 항목" in out
    assert "두 번째 항목" in out
    assert "| 구분 | 금액 |" in out


def test_parse_docx_extracts_paragraph_text(tmp_path):
    path = str(tmp_path / "doc.docx")
    _make_docx(
        path,
        "<w:p><w:r><w:t>안녕하세요</w:t></w:r></w:p>"
        "<w:p><w:r><w:t>두 </w:t></w:r><w:r><w:t>번째 &amp; 줄</w:t></w:r></w:p>",
    )

    out = DocumentService().parse_to_markdown(path)

    assert "안녕하세요" in out
    # runs within a paragraph join; entities are unescaped
    assert "두 번째 & 줄" in out


def test_parse_csv_and_txt_are_read_as_text(tmp_path):
    csv_path = tmp_path / "data.csv"
    csv_path.write_text("a,b\n1,2\n", encoding="utf-8")
    assert DocumentService().parse_to_markdown(str(csv_path)) == "a,b\n1,2\n"


@pytest.mark.parametrize("extension", [".xlsx", ".xlsm"])
def test_parse_modern_excel_without_libreoffice(tmp_path, extension):
    path = str(tmp_path / f"data{extension}")
    _make_xlsx(path)

    out = DocumentService().parse_to_markdown(path)

    assert "## 요약" in out
    assert "| 항목 | 금액 |" in out
    assert "| 매출 | 1200 | TRUE |" in out
    assert "|  | =SUM(B2:B2) |" in out


def test_unsupported_format_raises_clear_error():
    with pytest.raises(DocumentConversionError):
        DocumentService().parse_to_markdown("legacy.xls")


def test_load_office_document_uses_native_fallback_without_soffice(tmp_path):
    loader = DocumentsLoader(file_paths=[])
    loader.document_conversion_service.is_soffice_available = MagicMock(
        return_value=False
    )
    loader.document_service = MagicMock()
    loader.document_service.parse_to_markdown.return_value = "native text"

    result = loader.load_office_document("/tmp/foo.pptx", temp_dir=None)

    assert result == "native text"
    loader.document_service.parse_to_markdown.assert_called_once_with("/tmp/foo.pptx")


def test_load_office_document_parses_xlsx_without_soffice(tmp_path):
    path = str(tmp_path / "spreadsheet.xlsx")
    _make_xlsx(path)
    loader = DocumentsLoader(file_paths=[])
    loader.document_conversion_service.is_soffice_available = MagicMock(
        return_value=False
    )

    result = loader.load_office_document(path, temp_dir=None)

    assert "1200" in result
    assert "TRUE" in result
    assert "SUM(B2:B2)" in result


def test_load_office_document_errors_when_no_soffice_and_no_fallback():
    loader = DocumentsLoader(file_paths=[])
    loader.document_conversion_service.is_soffice_available = MagicMock(
        return_value=False
    )
    loader.document_service = None

    with pytest.raises(DocumentConversionError):
        loader.load_office_document("/tmp/foo.pptx", temp_dir=None)
