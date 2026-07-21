"""Unit tests for the authored generation pipeline (no real LLM / no real Chrome).

Covers the resilience guarantees that keep authored generation from aborting on one
bad slide: HTML validity gating, branded fallback, per-slide author retry/fallback,
per-slide render placeholder, style-brief propagation, image-PPTX assembly, the
authored-deck predicate, and brand resolution."""

import asyncio
import hashlib
import io
import uuid
from types import SimpleNamespace

from PIL import Image

from models.generate_presentation_request import GeneratePresentationRequest
from models.presentation_outline_model import (
    PresentationOutlineModel,
    SlideOutlineModel,
)
from models.sql.presentation import PresentationModel
from services import authored_presentation_service as authored_service
from services.authored_presentation_service import resolve_brand
from utils import slide_capture
from utils.llm_calls import author_deck as author_deck_mod
from utils.llm_calls import author_slide as author_slide_mod
from utils.llm_calls import author_vision_qa as author_vision_qa_mod
from utils.llm_calls.author_deck import (
    AuthoredDeckResult,
    author_deck,
    build_image_pptx,
    plan_deck_roles,
)
from utils.llm_calls.author_slide import (
    Brand,
    apply_style_defaults,
    build_design_system,
    fallback_slide_html,
    is_valid_slide_html,
)
from utils.llm_calls.critique_slide import SlideCritique


def _run(coro):
    return asyncio.run(coro)


_VALID = (
    "<!DOCTYPE html><html lang='ko'><head><meta charset='utf-8'></head>"
    "<body><h1>제목</h1></body></html>"
)


def _outline(n):
    return PresentationOutlineModel(
        slides=[SlideOutlineModel(content=f"슬라이드 {i} 내용") for i in range(n)]
    )


def _brand():
    return Brand(topic="주제", language="Korean", primary="#2563EB", fonts="Noto Sans KR")


# --- HTML validity + fallback ------------------------------------------------


def test_is_valid_slide_html_accepts_full_document():
    assert is_valid_slide_html(_VALID) is True


def test_is_valid_slide_html_rejects_garbage():
    assert is_valid_slide_html("") is False
    assert is_valid_slide_html("   ") is False
    assert is_valid_slide_html("not html") is False
    assert is_valid_slide_html("<html><head></head></html>") is False  # no body


def test_fallback_slide_html_is_valid_and_contains_content():
    html = fallback_slide_html("핵심 메시지 텍스트", _brand(), "COVER", 0, 5)
    assert is_valid_slide_html(html)
    assert "핵심 메시지 텍스트" in html
    assert "COVER" in html
    assert "#2563EB" in html


def test_fallback_slide_html_escapes_angle_brackets():
    html = fallback_slide_html("<script>x</script>", _brand(), "PROBLEM", 1, 3)
    assert "<script>" not in html
    assert "&lt;script&gt;" in html


# --- author_deck resilience (mocked LLM) -------------------------------------


def test_author_deck_uses_authored_html_when_valid(monkeypatch):
    async def fake(content, ds, brand, role, index, n):
        return _VALID

    monkeypatch.setattr(author_deck_mod, "author_slide_html", fake)
    htmls = _run(author_deck(_outline(4), _brand())).htmls
    assert len(htmls) == 4
    assert all(h == _VALID for h in htmls)


def test_author_deck_falls_back_when_authoring_raises(monkeypatch):
    async def boom(content, ds, brand, role, index, n):
        raise RuntimeError("provider down")

    monkeypatch.setattr(author_deck_mod, "author_slide_html", boom)
    htmls = _run(author_deck(_outline(3), _brand())).htmls
    assert len(htmls) == 3
    # Every slide degraded to a valid branded fallback (deck still completes).
    assert all(is_valid_slide_html(h) for h in htmls)
    assert any("슬라이드 0 내용" in h for h in htmls)


def test_author_deck_falls_back_when_authoring_returns_invalid(monkeypatch):
    async def empty(content, ds, brand, role, index, n):
        return ""

    monkeypatch.setattr(author_deck_mod, "author_slide_html", empty)
    htmls = _run(author_deck(_outline(2), _brand())).htmls
    assert all(is_valid_slide_html(h) for h in htmls)


def test_author_deck_retries_then_succeeds(monkeypatch):
    calls = {"n": 0}

    async def flaky(content, ds, brand, role, index, n):
        calls["n"] += 1
        if calls["n"] == 1:
            return ""  # first attempt invalid -> retry
        return _VALID

    monkeypatch.setattr(author_deck_mod, "author_slide_html", flaky)
    htmls = _run(author_deck(_outline(1), _brand())).htmls
    assert htmls == [_VALID]
    assert calls["n"] == 2


def test_default_design_system_is_byte_for_byte_legacy_output():
    expected_sha256 = "5ffac30c247a543fef5faa8b3a1c9dc945cc658f65fac243e269a5694e4f5c28"
    without_style = build_design_system(_brand())
    canonical_default = build_design_system(
        _brand(), SimpleNamespace(id="default", brief="this must not replace legacy")
    )

    assert len(without_style) == 2616
    assert hashlib.sha256(without_style.encode()).hexdigest() == expected_sha256
    assert canonical_default == without_style


def test_custom_design_system_injects_brief_without_weakening_common_rules():
    style = SimpleNamespace(
        id="editorial",
        brief="  - Editorial composition with raw {curly} braces.\n- Restrained typography.  ",
    )
    design_system = build_design_system(_brand(), style)

    assert "- Editorial composition with raw {curly} braces." in design_system
    assert "Brand primary colour: #2563EB" in design_system
    assert "EXACTLY 1280x720px" in design_system
    assert "FIT IS CRITICAL" in design_system
    assert "Return ONLY the complete HTML document." in design_system
    assert "The primary is the ONE accent" not in design_system


def test_author_deck_builds_and_shares_selected_design_system_once(monkeypatch):
    style = SimpleNamespace(id="editorial", brief="- Editorial")
    build_calls = []
    author_design_systems = []

    def fake_build(brand, selected_style):
        build_calls.append((brand, selected_style))
        return "SELECTED DESIGN SYSTEM"

    async def fake_author(content, design_system, brand, role, index, n):
        author_design_systems.append(design_system)
        return _VALID

    monkeypatch.setattr(author_deck_mod, "build_design_system", fake_build)
    monkeypatch.setattr(author_deck_mod, "author_slide_html", fake_author)

    result = _run(author_deck(_outline(3), _brand(), style))

    assert len(build_calls) == 1
    assert build_calls[0][1] is style
    assert result == AuthoredDeckResult(
        htmls=[_VALID, _VALID, _VALID], design_system="SELECTED DESIGN SYSTEM"
    )
    assert author_design_systems == ["SELECTED DESIGN SYSTEM"] * 3


def test_author_deck_sends_role_specific_visual_references(monkeypatch, tmp_path):
    cover = tmp_path / "cover.png"
    content = tmp_path / "content.png"
    closing = tmp_path / "closing.png"
    style = SimpleNamespace(
        id="visual",
        brief="- Visual",
        reference_images={
            "cover": (cover,),
            "content": (content,),
            "closing": (closing,),
        },
    )
    captured = []

    async def fake_author(
        content_text,
        design_system,
        brand,
        role,
        index,
        n,
        reference_images=None,
    ):
        captured.append((role, tuple(reference_images or ())))
        return _VALID

    monkeypatch.setattr(author_deck_mod, "author_slide_html", fake_author)

    _run(author_deck(_outline(3), _brand(), style))

    assert captured == [
        ("COVER", (cover,)),
        ("PROBLEM", (content,)),
        ("CLOSING", (closing,)),
    ]


def test_style_defaults_preserve_explicit_brand_overrides():
    style = SimpleNamespace(
        id="visual",
        primary_color="#35F2C2",
        body_font="Noto Sans KR",
        heading_font="Space Grotesk",
    )
    preset_brand = apply_style_defaults(_brand(), style)
    custom_brand = apply_style_defaults(
        Brand(
            topic="topic",
            primary="#FF6600",
            fonts="Pretendard",
            primary_is_explicit=True,
            fonts_are_explicit=True,
        ),
        style,
    )

    assert preset_brand.primary == "#35F2C2"
    assert preset_brand.fonts == "Noto Sans KR"
    assert custom_brand.primary == "#FF6600"
    assert custom_brand.fonts == "Pretendard"


def test_style_defaults_fall_back_to_preview_accent_without_a_primary_token():
    # 28 of 30 presets describe their palette in prose and carry no primary_color
    # token; the picker swatch accent must then become the effective brand primary
    # so the prompt and persisted theme match the style instead of the brand blue.
    style = SimpleNamespace(id="exec-report", preview_accent="#00A3E0")
    assert apply_style_defaults(_brand(), style).primary == "#00A3E0"

    # An explicit request colour still wins over the swatch fallback.
    explicit = Brand(topic="t", primary="#FF6600", primary_is_explicit=True)
    assert apply_style_defaults(explicit, style).primary == "#FF6600"

    # A structured primary_color token still takes precedence over the swatch.
    tokened = SimpleNamespace(
        id="cyber", preview_accent="#00A3E0", primary_color="#35F2C2"
    )
    assert apply_style_defaults(_brand(), tokened).primary == "#35F2C2"


def test_generate_request_accepts_optional_authored_style():
    assert GeneratePresentationRequest(content="topic").authored_style is None
    request = GeneratePresentationRequest(content="topic", authored_style="editorial")
    assert request.authored_style == "editorial"
    assert request.model_dump()["authored_style"] == "editorial"


def test_vision_qa_reuses_selected_design_system(monkeypatch):
    captured = []

    async def fake_critique(pngs, contexts=None):
        return [SlideCritique(needs_fix=True)]

    async def fake_author(
        content, design_system, brand, role, index, n, reasoning_effort=None
    ):
        captured.append((design_system, reasoning_effort))
        return _VALID

    async def fake_render(html, timeout=60):
        return b"revised-png"

    monkeypatch.setattr(author_vision_qa_mod, "critique_authored", fake_critique)
    monkeypatch.setattr(author_vision_qa_mod, "author_slide_html", fake_author)
    monkeypatch.setattr(author_vision_qa_mod, "render_html_to_png", fake_render)

    htmls, pngs, fixed = _run(
        author_vision_qa_mod.revise_authored_deck(
            [_VALID],
            [b"original-png"],
            ["content"],
            ["COVER"],
            _brand(),
            "SELECTED DESIGN SYSTEM",
        )
    )

    assert captured == [("SELECTED DESIGN SYSTEM", "high")]
    assert htmls == [_VALID]
    assert pngs == [b"revised-png"]
    assert fixed == [0]


def test_authored_service_resolves_passes_and_persists_canonical_style(monkeypatch):
    selected_style = SimpleNamespace(id="editorial", brief="- Editorial")
    captured = {}

    async def fake_author_deck(outline, brand, style):
        captured["style"] = style
        return AuthoredDeckResult([_VALID], "SELECTED DESIGN SYSTEM")

    async def fake_render(htmls):
        return [b"png"]

    async def fake_revise(
        htmls, pngs, contents, roles, brand, design_system, max_cycles, style=None
    ):
        captured["vision_design_system"] = design_system
        captured["vision_style"] = style
        return htmls, pngs, []

    class FakeSession:
        def __init__(self):
            self.added = []

        def add(self, value):
            self.added.append(value)

        def add_all(self, values):
            self.added.extend(values)

        async def commit(self):
            return None

    monkeypatch.setattr(
        authored_service,
        "resolve_authored_style",
        lambda style_id: selected_style if style_id == "editorial-alias" else None,
    )
    monkeypatch.setattr(authored_service, "author_deck", fake_author_deck)
    monkeypatch.setattr(authored_service, "render_html_list_to_pngs", fake_render)
    monkeypatch.setattr(authored_service, "revise_authored_deck", fake_revise)
    monkeypatch.setattr(authored_service, "find_chrome", lambda: True)
    monkeypatch.setattr(
        authored_service, "_save_slide_pngs", lambda presentation_id, pngs: ["slide.png"]
    )
    monkeypatch.setattr(
        authored_service,
        "_build_authored_export",
        lambda request, presentation_id, pngs, title: "deck.pptx",
    )

    from utils import llm_provider

    monkeypatch.setattr(
        llm_provider, "get_llm_provider", lambda: SimpleNamespace(value="stub")
    )
    monkeypatch.setattr(llm_provider, "get_model", lambda: "stub-model")

    session = FakeSession()
    request = GeneratePresentationRequest(
        content="topic",
        template="authored",
        authored_style="editorial-alias",
        vision_qa=True,
    )
    result = _run(
        authored_service.generate_authored_presentation(
            request, uuid.uuid4(), _outline(1), "English", session
        )
    )
    presentation = next(x for x in session.added if isinstance(x, PresentationModel))

    assert captured["style"] is selected_style
    assert captured["vision_design_system"] == "SELECTED DESIGN SYSTEM"
    assert captured["vision_style"] is selected_style
    assert presentation.theme["style"] == "editorial"
    assert presentation.theme["mode"] == "authored"
    assert result.path == "deck.pptx"


# --- render resilience (mocked Chrome) ---------------------------------------


def test_render_list_uses_placeholder_on_failure(monkeypatch):
    async def boom(html, timeout=60):
        raise RuntimeError("no chrome")

    monkeypatch.setattr(slide_capture, "render_html_to_png", boom)
    pngs = _run(slide_capture.render_html_list_to_pngs([_VALID, _VALID]))
    assert len(pngs) == 2
    for png in pngs:
        img = Image.open(io.BytesIO(png))
        assert img.size == (slide_capture.SLIDE_W, slide_capture.SLIDE_H)


# --- image PPTX assembly -----------------------------------------------------


def _png_1280x720(color):
    img = Image.new("RGB", (1280, 720), color)
    buf = io.BytesIO()
    img.save(buf, format="PNG")
    return buf.getvalue()


def test_build_image_pptx_one_fullbleed_picture_per_slide(tmp_path):
    from pptx import Presentation

    pngs = [_png_1280x720((i * 10, 0, 0)) for i in range(3)]
    out = str(tmp_path / "deck.pptx")
    build_image_pptx(pngs, out)
    prs = Presentation(out)
    assert len(prs.slides._sldIdLst) == 3
    assert prs.slide_width == 12191695 and prs.slide_height == 6858000  # 13.333x7.5"
    for s in prs.slides:
        pics = [sh for sh in s.shapes if sh.shape_type == 13]
        assert len(pics) == 1


def test_build_image_pdf_is_valid_multipage(tmp_path):
    from services.authored_presentation_service import _build_image_pdf

    pngs = [_png_1280x720((0, i * 20, 0)) for i in range(4)]
    out = str(tmp_path / "deck.pdf")
    _build_image_pdf(pngs, out)
    with open(out, "rb") as f:
        data = f.read()
    assert data[:5] == b"%PDF-"  # valid PDF magic (PIL can write but not read PDFs)
    assert b"/Count 4" in data  # four pages, one per slide image
    assert len(data) > 4000


# --- authored predicate + roles + brand --------------------------------------


def test_is_authored_predicate():
    by_theme = PresentationModel(
        content="x", n_slides=1, language="ko", theme={"mode": "authored"},
        layout={"name": "ignored"},
    )
    by_null_layout = PresentationModel(content="x", n_slides=1, language="ko")
    templated = PresentationModel(
        content="x", n_slides=1, language="ko", layout={"name": "adaptive"},
    )
    assert by_theme.is_authored() is True
    assert by_null_layout.is_authored() is True
    assert templated.is_authored() is False


def test_mode_column_is_authoritative_over_sentinels():
    # An explicit mode wins over the legacy theme/layout sentinels.
    authored_with_layout = PresentationModel(
        content="x", n_slides=1, language="ko",
        mode="authored", layout={"name": "korean-biz"},
    )
    template_without_layout = PresentationModel(
        content="x", n_slides=1, language="ko",
        mode="template", layout=None,
    )
    assert authored_with_layout.is_authored() is True
    assert template_without_layout.is_authored() is False


def test_get_new_presentation_preserves_mode_and_theme():
    # /derive previously dropped theme (losing authored brand colours + the legacy
    # sentinel); mode + theme must carry to the copy.
    original = PresentationModel(
        content="x", n_slides=1, language="ko",
        mode="authored", theme={"mode": "authored", "primary": "#123456"},
    )
    clone = original.get_new_presentation()
    assert clone.id != original.id
    assert clone.mode == "authored"
    assert clone.theme == {"mode": "authored", "primary": "#123456"}
    assert clone.is_authored() is True


def test_plan_deck_roles_cover_and_closing():
    roles = plan_deck_roles(_outline(5))
    assert roles[0] == "COVER"
    assert roles[-1] == "CLOSING"


def test_resolve_brand_language_fonts():
    req = GeneratePresentationRequest(content="t", language="Korean (한국어)", template="authored")
    ko = resolve_brand(req, _outline(2), "Korean (한국어)")
    assert ko.fonts == "Noto Sans KR"
    assert ko.primary == "#2563EB"

    req_en = GeneratePresentationRequest(content="t", language="English", template="authored")
    en = resolve_brand(req_en, _outline(2), "English")
    assert en.fonts == "Inter"


def test_resolve_brand_honours_custom_request_fields():
    req = GeneratePresentationRequest(
        content="제조업 스마트팩토리 전환 전략 2026",
        language="Korean",
        template="authored",
        primary_color="#FF6600",
        fonts="Pretendard",
        wordmark="ACME",
    )
    brand = resolve_brand(req, _outline(3), "Korean")
    assert brand.primary == "#FF6600"
    assert brand.fonts == "Pretendard"
    assert brand.wordmark == "ACME"


def test_authored_title_prefers_clean_content_over_verbose_outline():
    from services.authored_presentation_service import authored_title

    outline = PresentationOutlineModel(
        slides=[
            SlideOutlineModel(
                content="# 표지\n발표자 [발표자 이름]  날짜 2026-06-08  개요 매우 긴 생성 텍스트"
            ),
            SlideOutlineModel(content="본문"),
        ]
    )
    req = GeneratePresentationRequest(
        content="제조업 스마트팩토리 전환 전략 2026", template="authored"
    )
    title = authored_title(req, outline)
    assert title == "제조업 스마트팩토리 전환 전략 2026"
    assert "발표자" not in title
