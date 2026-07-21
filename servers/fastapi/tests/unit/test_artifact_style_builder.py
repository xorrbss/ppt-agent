from __future__ import annotations

import base64
import hashlib
import json
import re
from pathlib import Path
import subprocess
import sys
from typing import Iterable
from zipfile import ZIP_DEFLATED, ZipFile
import zlib

from PIL import Image
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.util import Inches, Pt
import pytest

from utils import artifact_style_analysis
from utils import artifact_style_pdf
from utils.artifact_style_analysis import (
    ArtifactAnalysisError,
    analysis_json_bytes,
    analyze_artifact,
)
from utils.artifact_style_builder import build_authored_style
from utils.authored_style_converter import ConversionError
from utils.authored_styles import AUTHORED_STYLES_DIRECTORY, load_authored_styles


REPOSITORY_ROOT = Path(__file__).resolve().parents[4]
CLI = REPOSITORY_ROOT / "scripts" / "build-authored-style.py"


def _write_pdf(path: Path, objects: Iterable[bytes]) -> None:
    payload = bytearray(b"%PDF-1.4\n%\xe2\xe3\xcf\xd3\n")
    offsets = [0]
    object_list = list(objects)
    for number, body in enumerate(object_list, start=1):
        offsets.append(len(payload))
        payload.extend(f"{number} 0 obj\n".encode())
        payload.extend(body)
        payload.extend(b"\nendobj\n")
    xref = len(payload)
    payload.extend(f"xref\n0 {len(object_list) + 1}\n".encode())
    payload.extend(b"0000000000 65535 f \n")
    for offset in offsets[1:]:
        payload.extend(f"{offset:010d} 00000 n \n".encode())
    payload.extend(
        (
            f"trailer\n<< /Size {len(object_list) + 1} /Root 1 0 R >>\n"
            f"startxref\n{xref}\n%%EOF\n"
        ).encode()
    )
    path.write_bytes(bytes(payload))


def _stream(data: bytes, extra: bytes = b"") -> bytes:
    return (
        b"<< /Length "
        + str(len(data)).encode()
        + b" "
        + extra
        + b">>\nstream\n"
        + data
        + b"\nendstream"
    )


def _text_pdf(path: Path) -> None:
    content = (
        b"BT /F1 24 Tf 72 700 Td 0.10 0.20 0.80 rg (Evidence Title) Tj ET\n"
        b"BT /F1 12 Tf 72 650 Td 0.05 0.05 0.05 rg (Body text) Tj ET\n"
        b"0.10 0.20 0.80 rg 72 500 200 80 re f\n"
    )
    _write_pdf(
        path,
        [
            b"<< /Type /Catalog /Pages 2 0 R >>",
            b"<< /Type /Pages /Kids [3 0 R 4 0 R] /Count 2 >>",
            b"<< /Type /Page /Parent 2 0 R /MediaBox [0 0 612 792] /Resources << /Font << /F1 5 0 R >> >> /Contents 6 0 R >>",
            b"<< /Type /Page /Parent 2 0 R /MediaBox [0 0 612 792] /Resources << /Font << /F1 5 0 R >> >> /Contents 7 0 R >>",
            b"<< /Type /Font /Subtype /Type1 /BaseFont /Helvetica >>",
            _stream(content),
            _stream(content),
        ],
    )


def _image_pdf(path: Path) -> None:
    image_data = b"FF0000>"
    image = _stream(
        image_data,
        b"/Type /XObject /Subtype /Image /Width 1 /Height 1 "
        b"/ColorSpace /DeviceRGB /BitsPerComponent 8 /Filter /ASCIIHexDecode ",
    )
    content = b"q 612 0 0 792 0 0 cm /Im0 Do Q\n"
    _write_pdf(
        path,
        [
            b"<< /Type /Catalog /Pages 2 0 R >>",
            b"<< /Type /Pages /Kids [3 0 R] /Count 1 >>",
            b"<< /Type /Page /Parent 2 0 R /MediaBox [0 0 612 792] /Resources << /XObject << /Im0 4 0 R >> >> /Contents 5 0 R >>",
            image,
            _stream(content),
        ],
    )


def _empty_pdf(path: Path) -> None:
    _write_pdf(
        path,
        [
            b"<< /Type /Catalog /Pages 2 0 R >>",
            b"<< /Type /Pages /Kids [] /Count 0 >>",
        ],
    )


def _blank_page_pdf(
    path: Path, *, width: int = 612, height: int = 792, catalog_extra: bytes = b""
) -> None:
    _write_pdf(
        path,
        [
            b"<< /Type /Catalog /Pages 2 0 R " + catalog_extra + b">>",
            b"<< /Type /Pages /Kids [3 0 R] /Count 1 >>",
            (
                b"<< /Type /Page /Parent 2 0 R /MediaBox [0 0 "
                + str(width).encode()
                + b" "
                + str(height).encode()
                + b"] /Contents 4 0 R >>"
            ),
            _stream(b""),
        ],
    )


def _active_content_pdf(path: Path) -> None:
    _write_pdf(
        path,
        [
            b"<< /Type /Catalog /Pages 2 0 R /OpenAction 5 0 R "
            b"/Names << /EmbeddedFiles 6 0 R >> >>",
            b"<< /Type /Pages /Kids [3 0 R] /Count 1 >>",
            b"<< /Type /Page /Parent 2 0 R /MediaBox [0 0 612 792] /Contents 4 0 R >>",
            _stream(b""),
            b"<< /S /JavaScript /JS (inert-test-marker) >>",
            b"<< /Names [(payload.bin) 7 0 R] >>",
            b"<< /Type /Filespec /F (payload.bin) /EF << /F 8 0 R >> >>",
            _stream(b"", b"/Type /EmbeddedFile "),
        ],
    )


def _flate_pdf(path: Path, expanded: bytes) -> None:
    compressed = zlib.compress(expanded, level=9)
    _write_pdf(
        path,
        [
            b"<< /Type /Catalog /Pages 2 0 R >>",
            b"<< /Type /Pages /Kids [3 0 R] /Count 1 >>",
            b"<< /Type /Page /Parent 2 0 R /MediaBox [0 0 612 792] /Contents 4 0 R >>",
            _stream(compressed, b"/Filter /FlateDecode "),
        ],
    )


def _filter_chain_pdf(path: Path) -> None:
    content = b"BT /F1 18 Tf 72 700 Td (Bounded filter chain) Tj ET\n"
    encoded = base64.a85encode(zlib.compress(content), adobe=True)
    _write_pdf(
        path,
        [
            b"<< /Type /Catalog /Pages 2 0 R >>",
            b"<< /Type /Pages /Kids [3 0 R] /Count 1 >>",
            b"<< /Type /Page /Parent 2 0 R /MediaBox [0 0 612 792] /Resources << /Font << /F1 5 0 R >> >> /Contents 4 0 R >>",
            _stream(encoded, b"/Filter [/ASCII85Decode /FlateDecode] "),
            b"<< /Type /Font /Subtype /Type1 /BaseFont /Helvetica >>",
        ],
    )


def _run_length_pdf(path: Path, expanded_size: int) -> None:
    encoded = bytearray()
    remaining = expanded_size
    while remaining:
        count = min(128, remaining)
        encoded.extend((257 - count, ord("A")))
        remaining -= count
    encoded.append(128)
    _write_pdf(
        path,
        [
            b"<< /Type /Catalog /Pages 2 0 R >>",
            b"<< /Type /Pages /Kids [3 0 R] /Count 1 >>",
            b"<< /Type /Page /Parent 2 0 R /MediaBox [0 0 612 792] /Contents 4 0 R >>",
            _stream(bytes(encoded), b"/Filter /RunLengthDecode "),
        ],
    )


def _escaped_and_compressed_active_content_pdf(path: Path) -> None:
    object_stream = zlib.compress(b"9 0 << /S /Launch /F (never-opened.bin) >>")
    _write_pdf(
        path,
        [
            b"<< /Type /Catalog /Pages 2 0 R /Open#41ction 5 0 R >>",
            b"<< /Type /Pages /Kids [3 0 R] /Count 1 >>",
            b"<< /Type /Page /Parent 2 0 R /MediaBox [0 0 612 792] /Contents 4 0 R >>",
            _stream(b""),
            b"<< /S /Java#53cript /J#53 (inert-test-marker) >>",
            _stream(
                object_stream,
                b"/Type /ObjStm /N 1 /First 4 /Filter /FlateDecode ",
            ),
        ],
    )


def _pptx(path: Path, image_path: Path) -> None:
    Image.new("RGB", (16, 16), "#22AA88").save(image_path)
    presentation = Presentation()
    presentation.slide_width = Inches(13.333)
    presentation.slide_height = Inches(7.5)
    presentation.core_properties.title = "Deterministic fixture"
    for index in range(2):
        slide = presentation.slides.add_slide(presentation.slide_layouts[6])
        title = slide.shapes.add_textbox(
            Inches(0.8), Inches(0.6), Inches(7), Inches(0.8)
        )
        run = title.text_frame.paragraphs[0].add_run()
        run.text = f"Observed title {index + 1}"
        run.font.name = "Aptos Display"
        run.font.size = Pt(28)
        run.font.color.rgb = RGBColor(15, 23, 42)
        panel = slide.shapes.add_shape(
            1, Inches(0.8), Inches(1.8), Inches(7), Inches(4.6)
        )
        panel.fill.solid()
        panel.fill.fore_color.rgb = RGBColor(241, 245, 249)
        panel.line.color.rgb = RGBColor(37, 99, 235)
        slide.shapes.add_picture(
            str(image_path), Inches(9.2), Inches(2), Inches(2.5), Inches(2.5)
        )
    presentation.save(path)


def _add_inert_security_parts(path: Path) -> None:
    relationships = (
        b'<?xml version="1.0" encoding="UTF-8"?>'
        b'<Relationships xmlns="http://schemas.openxmlformats.org/package/2006/relationships">'
        b'<Relationship Id="rId1" Type="urn:test" Target="https://example.invalid/" '
        b'TargetMode="External"/></Relationships>'
    )
    with ZipFile(path, "a", compression=ZIP_DEFLATED) as package:
        package.writestr("ppt/vbaProject.bin", b"inert-test-marker")
        package.writestr("ppt/embeddings/object.bin", b"inert-test-marker")
        package.writestr("ppt/activeX/control.bin", b"inert-test-marker")
        package.writestr("custom/_rels/security.rels", relationships)


def _damage_zip_member(path: Path, member_name: str) -> None:
    with ZipFile(path) as package:
        member = package.getinfo(member_name)
    raw = bytearray(path.read_bytes())
    offset = member.header_offset
    filename_length = int.from_bytes(raw[offset + 26 : offset + 28], "little")
    extra_length = int.from_bytes(raw[offset + 28 : offset + 30], "little")
    payload_offset = offset + 30 + filename_length + extra_length
    raw[payload_offset + max(0, member.compress_size // 2)] ^= 0x01
    path.write_bytes(raw)


def _catalog_digest() -> str:
    digest = hashlib.sha256()
    for source in sorted(AUTHORED_STYLES_DIRECTORY.glob("*.yaml")):
        digest.update(source.name.encode())
        digest.update(source.read_bytes())
    return digest.hexdigest()


def _warning_codes(analysis: dict) -> set[str]:
    return {warning["code"] for warning in analysis["warnings"]}


def _rewrite_presentation_xml(path: Path, transform) -> None:
    with ZipFile(path) as package:
        members = {name: package.read(name) for name in package.namelist()}
    xml = members["ppt/presentation.xml"].decode("utf-8")
    members["ppt/presentation.xml"] = transform(xml).encode("utf-8")
    with ZipFile(path, "w", compression=ZIP_DEFLATED) as package:
        for name, data in members.items():
            package.writestr(name, data)


def test_pptx_analysis_fails_clean_on_missing_or_zero_slide_size(
    tmp_path: Path,
) -> None:
    # A malformed presentation.xml (no <p:sldSz>, or a zero dimension) must raise
    # ArtifactAnalysisError — caught by the CLI's clean handler — rather than a raw
    # TypeError/ZeroDivisionError traceback from the aspect-ratio math.
    missing = tmp_path / "no-size.pptx"
    _pptx(missing, tmp_path / "pixel.png")
    _rewrite_presentation_xml(missing, lambda xml: re.sub(r"<p:sldSz[^>]*/>", "", xml))
    with pytest.raises(ArtifactAnalysisError, match="slide size"):
        analyze_artifact(missing)

    zero = tmp_path / "zero-size.pptx"
    _pptx(zero, tmp_path / "pixel2.png")
    _rewrite_presentation_xml(
        zero, lambda xml: re.sub(r'cy="\d+"', 'cy="0"', xml, count=1)
    )
    with pytest.raises(ArtifactAnalysisError, match="invalid slide dimensions"):
        analyze_artifact(zero)


def test_pptx_analysis_is_deterministic_and_extracts_design_signals(
    tmp_path: Path,
) -> None:
    source = tmp_path / "observed-deck.pptx"
    _pptx(source, tmp_path / "pixel.png")

    first = analyze_artifact(source)
    second = analyze_artifact(source)

    assert analysis_json_bytes(first) == analysis_json_bytes(second)
    assert first["document"]["page_count"] == 2
    assert first["signals"]["colors"]["status"] == "observed"
    assert first["signals"]["fonts"]["values"][0]["family"] == "Aptos Display"
    assert first["signals"]["repeated_layouts"]["values"][0]["count"] == 2
    assert first["signals"]["composition"]["element_counts"]["image"] == 2
    assert first["security"]["external_relationship_count"] == 0


def test_pdf_analysis_is_deterministic_and_strips_subset_independent_fonts(
    tmp_path: Path,
) -> None:
    source = tmp_path / "observed-pages.pdf"
    _text_pdf(source)

    first = analyze_artifact(source)
    second = analyze_artifact(source)

    assert analysis_json_bytes(first) == analysis_json_bytes(second)
    assert first["document"]["page_count"] == 2
    assert first["document"]["page_size"]["unit"] == "points"
    assert first["signals"]["fonts"]["values"][0]["family"] == "Helvetica"
    assert first["signals"]["text_hierarchy"]["status"] == "observed"
    assert first["signals"]["repeated_layouts"]["values"][0]["count"] == 2


def test_image_only_pdf_reports_scan_likelihood_without_ocr(tmp_path: Path) -> None:
    source = tmp_path / "scan.pdf"
    _image_pdf(source)

    analysis = analyze_artifact(source)

    assert "likely_scanned_pdf" in _warning_codes(analysis)
    assert "colors_unavailable" in _warning_codes(analysis)
    assert "text_hierarchy_unavailable" in _warning_codes(analysis)
    assert analysis["signals"]["fonts"]["confidence"] == "none"
    assert analysis["signals"]["composition"]["element_counts"]["image"] == 1


def test_pptx_analysis_reports_active_and_external_content(tmp_path: Path) -> None:
    source = tmp_path / "security-signals.pptx"
    _pptx(source, tmp_path / "pixel.png")
    _add_inert_security_parts(source)

    analysis = analyze_artifact(source)

    assert analysis["document"]["page_count"] == 2
    assert analysis["security"]["macro_part_count"] == 1
    assert analysis["security"]["embedded_part_count"] == 1
    assert analysis["security"]["active_content_part_count"] == 1
    assert analysis["security"]["external_relationship_count"] == 1
    assert {
        "pptx_macros_present",
        "pptx_embedded_content_present",
        "pptx_active_content_present",
        "pptx_external_links_present",
    }.issubset(_warning_codes(analysis))


def test_pdf_analysis_reports_active_and_embedded_content(tmp_path: Path) -> None:
    source = tmp_path / "security-signals.pdf"
    _active_content_pdf(source)

    analysis = analyze_artifact(source)

    assert analysis["security"]["active_content_counts"]["javascript"] == 2
    assert analysis["security"]["active_content_counts"]["open_actions"] == 1
    assert analysis["security"]["embedded_content_marker_count"] >= 1
    assert {
        "pdf_active_content_present",
        "pdf_embedded_content_present",
    }.issubset(_warning_codes(analysis))


def test_pdf_compressed_stream_expansion_is_bounded(
    tmp_path: Path, monkeypatch: pytest.MonkeyPatch
) -> None:
    source = tmp_path / "compressed-bomb.pdf"
    _flate_pdf(source, b"A" * 16_384)
    monkeypatch.setattr(artifact_style_pdf, "MAX_PDF_EXPANDED_STREAM_BYTES", 1024)

    with pytest.raises(ArtifactAnalysisError, match="expanded-size safety limit"):
        analyze_artifact(source)


def test_pdf_stream_comment_cannot_bypass_preflight(
    tmp_path: Path, monkeypatch: pytest.MonkeyPatch
) -> None:
    source = tmp_path / "comment-before-stream.pdf"
    compressed = zlib.compress(b"A" * 8_192, level=9)
    stream = (
        b"<< /Length "
        + str(len(compressed)).encode()
        + b" /Filter /FlateDecode >> % lexical whitespace\nstream\n"
        + compressed
        + b"\nendstream"
    )
    _write_pdf(
        source,
        [
            b"<< /Type /Catalog /Pages 2 0 R >>",
            b"<< /Type /Pages /Kids [3 0 R] /Count 1 >>",
            b"<< /Type /Page /Parent 2 0 R /MediaBox [0 0 612 792] /Contents 4 0 R >>",
            stream,
        ],
    )
    monkeypatch.setattr(artifact_style_pdf, "MAX_PDF_EXPANDED_STREAM_BYTES", 1024)

    def _must_not_open(*args: object, **kwargs: object) -> None:
        pytest.fail("pdfplumber.open must not run before stream preflight")

    monkeypatch.setattr(artifact_style_pdf.pdfplumber, "open", _must_not_open)
    with pytest.raises(ArtifactAnalysisError, match="expanded-size safety limit"):
        analyze_artifact(source)


def test_pdf_unbounded_stream_dictionary_is_rejected_before_open(
    tmp_path: Path, monkeypatch: pytest.MonkeyPatch
) -> None:
    source = tmp_path / "oversized-stream-dictionary.pdf"
    _write_pdf(
        source,
        [
            b"<< /Type /Catalog /Pages 2 0 R >>",
            b"<< /Type /Pages /Kids [3 0 R] /Count 1 >>",
            b"<< /Type /Page /Parent 2 0 R /MediaBox [0 0 612 792] /Contents 4 0 R >>",
            _stream(b"", b"/Padding (" + b"A" * 256 + b") "),
        ],
    )
    monkeypatch.setattr(artifact_style_pdf, "MAX_PDF_STREAM_DICTIONARY_BYTES", 64)

    def _must_not_open(*args: object, **kwargs: object) -> None:
        pytest.fail("pdfplumber.open must not run before stream preflight")

    monkeypatch.setattr(artifact_style_pdf.pdfplumber, "open", _must_not_open)
    with pytest.raises(ArtifactAnalysisError, match="cannot be located"):
        analyze_artifact(source)


def test_pdf_run_length_expansion_is_bounded_before_open(
    tmp_path: Path, monkeypatch: pytest.MonkeyPatch
) -> None:
    source = tmp_path / "run-length-bomb.pdf"
    _run_length_pdf(source, 8_192)
    monkeypatch.setattr(artifact_style_pdf, "MAX_PDF_EXPANDED_STREAM_BYTES", 1024)

    def _must_not_open(*args: object, **kwargs: object) -> None:
        pytest.fail("pdfplumber.open must not run before stream preflight")

    monkeypatch.setattr(artifact_style_pdf.pdfplumber, "open", _must_not_open)
    with pytest.raises(ArtifactAnalysisError, match="expanded-size safety limit"):
        analyze_artifact(source)


def test_pdf_supported_filter_chain_is_decoded_in_order(tmp_path: Path) -> None:
    source = tmp_path / "filter-chain.pdf"
    _filter_chain_pdf(source)

    analysis = analyze_artifact(source)

    assert analysis["document"]["page_count"] == 1
    assert analysis["signals"]["fonts"]["values"][0]["family"] == "Helvetica"


def test_pdf_indirect_filter_is_rejected_before_open(
    tmp_path: Path, monkeypatch: pytest.MonkeyPatch
) -> None:
    source = tmp_path / "indirect-filter.pdf"
    compressed = zlib.compress(b"q Q\n")
    _write_pdf(
        source,
        [
            b"<< /Type /Catalog /Pages 2 0 R >>",
            b"<< /Type /Pages /Kids [3 0 R] /Count 1 >>",
            b"<< /Type /Page /Parent 2 0 R /MediaBox [0 0 612 792] /Contents 4 0 R >>",
            _stream(compressed, b"/Filter 5 0 R "),
            b"/FlateDecode",
        ],
    )

    def _must_not_open(*args: object, **kwargs: object) -> None:
        pytest.fail("pdfplumber.open must not run before stream preflight")

    monkeypatch.setattr(artifact_style_pdf.pdfplumber, "open", _must_not_open)
    with pytest.raises(ArtifactAnalysisError, match="indirect filter"):
        analyze_artifact(source)


def test_pdf_name_escapes_and_decoded_object_stream_markers_are_reported(
    tmp_path: Path,
) -> None:
    source = tmp_path / "escaped-security-signals.pdf"
    _escaped_and_compressed_active_content_pdf(source)

    analysis = analyze_artifact(source)

    assert analysis["security"]["active_content_counts"]["javascript"] == 2
    assert analysis["security"]["active_content_counts"]["open_actions"] == 1
    assert analysis["security"]["active_content_counts"]["launch_actions"] == 1
    assert "pdf_active_content_present" in _warning_codes(analysis)


def test_blank_pdf_reports_unavailable_composition(tmp_path: Path) -> None:
    source = tmp_path / "blank.pdf"
    _blank_page_pdf(source)

    analysis = analyze_artifact(source)

    assert analysis["signals"]["composition"]["status"] == "unavailable"
    assert analysis["signals"]["composition"]["confidence"] == "none"
    assert "pdf_design_objects_unavailable" in _warning_codes(analysis)


def test_build_writes_analysis_and_loader_validated_yaml_without_catalog_changes(
    tmp_path: Path,
) -> None:
    source = tmp_path / "observed-deck.pptx"
    target = tmp_path / "generated" / "observed.yaml"
    analysis_target = tmp_path / "generated" / "observed.analysis.json"
    _pptx(source, tmp_path / "pixel.png")
    before = _catalog_digest()

    result = build_authored_style(
        source,
        target,
        analysis_output=analysis_target,
        style_id="observed-deck-draft",
    )

    assert result.conversion.style_id == "observed-deck-draft"
    assert target.is_file()
    assert json.loads(analysis_target.read_text(encoding="utf-8")) == result.analysis
    loaded = load_authored_styles(target.parent)
    assert [style.id for style in loaded] == ["observed-deck-draft"]
    assert "Confidence:" in loaded[0].brief
    assert _catalog_digest() == before

    second_target = tmp_path / "generated" / "observed-second.yaml"
    build_authored_style(
        source,
        second_target,
        style_id="observed-deck-draft",
    )
    assert target.read_bytes() == second_target.read_bytes()


@pytest.mark.parametrize("artifact_format", ["pptx", "pdf"])
@pytest.mark.parametrize("dry_run", [True, False])
def test_build_fails_closed_on_active_content_without_writing(
    tmp_path: Path, artifact_format: str, dry_run: bool
) -> None:
    source = tmp_path / f"active-content.{artifact_format}"
    if artifact_format == "pptx":
        _pptx(source, tmp_path / "pixel.png")
        _add_inert_security_parts(source)
    else:
        _active_content_pdf(source)
    target = tmp_path / "draft.yaml"
    analysis_target = tmp_path / "draft.json"
    before = _catalog_digest()

    with pytest.raises(
        ArtifactAnalysisError, match="active or embedded content is not accepted"
    ):
        build_authored_style(
            source,
            target,
            analysis_output=analysis_target,
            dry_run=dry_run,
        )

    assert not target.exists()
    assert not analysis_target.exists()
    assert _catalog_digest() == before


@pytest.mark.parametrize("artifact_format", ["pptx", "pdf"])
def test_cli_fails_closed_on_active_content_without_writing(
    tmp_path: Path, artifact_format: str
) -> None:
    source = tmp_path / f"active-content.{artifact_format}"
    if artifact_format == "pptx":
        _pptx(source, tmp_path / "pixel.png")
        _add_inert_security_parts(source)
    else:
        _active_content_pdf(source)
    target = tmp_path / "draft.yaml"
    analysis_target = tmp_path / "draft.json"

    result = subprocess.run(
        [
            sys.executable,
            str(CLI),
            str(source),
            "-o",
            str(target),
            "--analysis-output",
            str(analysis_target),
        ],
        cwd=tmp_path,
        check=False,
        capture_output=True,
        text=True,
    )

    assert result.returncode == 2
    assert "active or embedded content is not accepted" in result.stderr
    assert "Traceback" not in result.stderr
    assert not target.exists()
    assert not analysis_target.exists()


def test_dry_run_writes_nothing_and_reuses_output_preflight(tmp_path: Path) -> None:
    source = tmp_path / "observed.pdf"
    target = tmp_path / "draft.yaml"
    analysis_target = tmp_path / "draft.json"
    _text_pdf(source)

    result = build_authored_style(
        source,
        target,
        analysis_output=analysis_target,
        dry_run=True,
    )

    assert result.conversion.target == target
    assert not target.exists()
    assert not analysis_target.exists()

    target.write_text("existing", encoding="utf-8")
    with pytest.raises(ConversionError, match="use --overwrite"):
        build_authored_style(source, target, dry_run=True)


@pytest.mark.parametrize("dry_run", [True, False])
def test_build_rejects_overlapping_output_paths_before_writing(
    tmp_path: Path, dry_run: bool
) -> None:
    source = tmp_path / "observed.pdf"
    target = tmp_path / "draft.yaml"
    _text_pdf(source)

    with pytest.raises(ConversionError, match="must not overlap"):
        build_authored_style(
            source,
            target,
            analysis_output=target / "analysis.json",
            dry_run=dry_run,
        )

    assert not target.exists()


def test_build_rejects_unsafe_id_and_invalid_analysis_suffix(tmp_path: Path) -> None:
    source = tmp_path / "observed.pdf"
    target = tmp_path / "draft.yaml"
    _text_pdf(source)

    with pytest.raises(ConversionError, match="unsafe style id"):
        build_authored_style(source, target, style_id="../unsafe")
    with pytest.raises(ConversionError, match="must use .json"):
        build_authored_style(source, target, analysis_output=target)


def test_encrypted_pdf_failure_is_clear(
    tmp_path: Path, monkeypatch: pytest.MonkeyPatch
) -> None:
    source = tmp_path / "encrypted.pdf"
    source.write_bytes(b"%PDF-1.4\n")

    def _password_failure(*args: object, **kwargs: object) -> None:
        raise artifact_style_pdf.PdfminerException(
            artifact_style_pdf.PDFPasswordIncorrect()
        )

    monkeypatch.setattr(artifact_style_pdf.pdfplumber, "open", _password_failure)
    with pytest.raises(ArtifactAnalysisError, match="requires a password"):
        analyze_artifact(source)


def test_clear_failures_for_empty_large_and_damaged_inputs(
    tmp_path: Path, monkeypatch: pytest.MonkeyPatch
) -> None:
    empty_deck = tmp_path / "empty.pptx"
    Presentation().save(empty_deck)
    with pytest.raises(ArtifactAnalysisError, match="contains no slides"):
        analyze_artifact(empty_deck)

    empty_pdf = tmp_path / "empty.pdf"
    _empty_pdf(empty_pdf)
    with pytest.raises(ArtifactAnalysisError, match="contains no pages"):
        analyze_artifact(empty_pdf)

    damaged = tmp_path / "damaged.pdf"
    damaged.write_bytes(b"%PDF-not-valid")
    with pytest.raises(ArtifactAnalysisError, match="damaged or unsupported PDF"):
        analyze_artifact(damaged)

    large = tmp_path / "large.pdf"
    large.write_bytes(b"%PDF-1.4")
    monkeypatch.setattr(artifact_style_analysis, "MAX_ARTIFACT_BYTES", 4)
    with pytest.raises(ArtifactAnalysisError, match="safety limit"):
        analyze_artifact(large)


def test_invalid_pdf_dimensions_and_damaged_pptx_member_are_clear(
    tmp_path: Path,
) -> None:
    invalid_dimensions = tmp_path / "invalid-dimensions.pdf"
    _blank_page_pdf(invalid_dimensions, height=0)
    with pytest.raises(ArtifactAnalysisError, match="invalid dimensions"):
        analyze_artifact(invalid_dimensions)

    damaged_pptx = tmp_path / "damaged-member.pptx"
    _pptx(damaged_pptx, tmp_path / "pixel.png")
    _damage_zip_member(damaged_pptx, "ppt/theme/theme1.xml")
    with pytest.raises(ArtifactAnalysisError, match="damaged or unsupported member"):
        analyze_artifact(damaged_pptx)

    cli_result = subprocess.run(
        [sys.executable, str(CLI), str(damaged_pptx), "-o", str(tmp_path / "x.yaml")],
        cwd=tmp_path,
        check=False,
        capture_output=True,
        text=True,
    )
    assert cli_result.returncode == 2
    assert "damaged or unsupported member" in cli_result.stderr
    assert "Traceback" not in cli_result.stderr


def test_cli_help_and_error_do_not_emit_tracebacks(tmp_path: Path) -> None:
    help_result = subprocess.run(
        [sys.executable, str(CLI), "--help"],
        cwd=tmp_path,
        check=False,
        capture_output=True,
        text=True,
    )
    assert help_result.returncode == 0
    assert "--analysis-output" in help_result.stdout
    assert "--dry-run" in help_result.stdout

    error_result = subprocess.run(
        [sys.executable, str(CLI), "missing.pdf", "-o", "draft.yaml"],
        cwd=tmp_path,
        check=False,
        capture_output=True,
        text=True,
    )
    assert error_result.returncode == 2
    assert "error:" in error_result.stderr
    assert "Traceback" not in error_result.stderr
