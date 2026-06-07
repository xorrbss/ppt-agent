"""Legacy-template editable-PPTX smoke (Docker/Linux/CI only).

Confirms a NON-adaptive (legacy `general` template) deck still exports to editable
PPTX with the current converter — a regression guard for converter-version bumps
(e.g. v0.2.9 -> v0.3.x). Same converter guard / skip behaviour as the adaptive G4
harness (scripts/check_adaptive_pptx_roundtrip.py).

Run (inside the Docker dev stack):
  RUN_PPTX_ROUNDTRIP=1 NEXT_PUBLIC_FAST_API=http://127.0.0.1:8000 \
    APP_DATA_DIRECTORY=/app_data uv run python scripts/check_legacy_pptx_roundtrip.py
"""
import asyncio
import os
import sys
import uuid

sys.path.insert(0, os.path.join(os.path.dirname(__file__), "..", "servers", "fastapi"))

from services.export_task_service import EXPORT_TASK_SERVICE  # noqa: E402


def _skip(reason: str) -> None:
    print(f"SKIP (legacy PPTX smoke): {reason}")
    sys.exit(0)


if not os.path.isfile(EXPORT_TASK_SERVICE.converter_path):
    _skip(
        f"export converter not found at {EXPORT_TASK_SERVICE.converter_path} "
        "- byte-level PPTX export needs the Linux/Docker runtime."
    )

from pptx import Presentation  # noqa: E402

from models.sql.presentation import PresentationModel  # noqa: E402
from models.sql.slide import SlideModel  # noqa: E402
from services.database import async_session_maker  # noqa: E402
from sqlalchemy import delete as sql_delete  # noqa: E402
from templates.get_layout_by_name import get_layout_by_name  # noqa: E402
from utils.export_utils import export_presentation  # noqa: E402

PID = uuid.UUID("1e6ac111-0000-4000-8000-000000000011")
GROUP = "korean-biz"
LAYOUT = "korean-biz:korean-biz-cover"  # registry layoutId is group-prefixed
# 1x1 transparent PNG so the legacy <img> has a valid src (no network fetch).
_PNG = (
    "data:image/png;base64,iVBORw0KGgoAAAANSUhEUgAAAAEAAAABCAQAAAC1HAwC"
    "AAAAC0lEQVR42mNkYAAAAAYAAjCB0C8AAAAASUVORK5CYII="
)
# Non-adaptive (concrete-TSX) regression guard: a korean-biz deck must still
# round-trip to editable PPTX via the converter. (Was the now-deleted `general`
# group; repointed to a kept curated group after the #5 legacy purge.)
SLIDES = [
    {
        "title": "비적응형 회귀 스모크",
        "subtitle": "한국형 비즈니스(비적응형) 템플릿 덱이 컨버터에서 편집가능 PPTX로 export되는지 확인",
        "presenterName": "테스트",
        "presentationDate": "2026년 6월",
        "image": {"__image_url__": _PNG, "__image_prompt__": "abstract cover"},
    },
    {
        "title": "두 번째 슬라이드",
        "subtitle": "멀티 슬라이드 비적응형 export 동작 확인 — 텍스트가 편집가능 텍스트프레임으로 나와야 함",
        "presenterName": "테스트",
        "presentationDate": "2026년 6월",
        "image": {"__image_url__": _PNG, "__image_prompt__": "abstract"},
    },
]


async def seed():
    layout = await get_layout_by_name(GROUP)
    idx = layout.get_slide_layout_index(LAYOUT)
    async with async_session_maker() as s:
        await s.execute(sql_delete(SlideModel).where(SlideModel.presentation == PID))
        old = await s.get(PresentationModel, PID)
        if old:
            await s.delete(old)
        await s.commit()
        s.add(PresentationModel(
            id=PID, content="legacy smoke", n_slides=len(SLIDES), language="Korean",
            title="legacy smoke", outlines={"slides": [{"content": "x"} for _ in SLIDES]},
            layout=layout.model_dump(),
            structure={"slides": [idx for _ in SLIDES]},
            tone="default", verbosity="standard",
        ))
        for i, c in enumerate(SLIDES):
            s.add(SlideModel(presentation=PID, layout_group=GROUP, layout=LAYOUT,
                             index=i, speaker_note="", content=c))
        await s.commit()


async def main():
    await seed()
    res = await export_presentation(PID, "legacy smoke", "pptx")
    prs = Presentation(res.path)
    print(f"exported {res.path} -> {len(prs.slides)} slides")
    failures: list = []

    def check(name, cond):
        print(f"  {'ok ' if cond else 'FAIL'} - {name}")
        if not cond:
            failures.append(name)

    all_text = " ".join(
        sh.text_frame.text for sl in prs.slides for sh in sl.shapes if sh.has_text_frame
    )
    check("slide count == n", len(prs.slides) == len(SLIDES))
    check("legacy title text editable", "비적응형 회귀 스모크" in all_text)
    check("legacy multi-slide text present", "두 번째 슬라이드" in all_text)

    if failures:
        print(f"\nLEGACY SMOKE FAIL: {failures}")
        sys.exit(1)
    print("\nLEGACY SMOKE PASS: legacy deck round-trips to editable PPTX.")


asyncio.run(main())
